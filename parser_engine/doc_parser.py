"""文档解析：.docx / .pdf / .pptx / .txt / .xlsx"""
import os


def parse_docx(file_path):
    from docx import Document
    doc = Document(file_path)
    paras = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paras)


def parse_pdf(file_path):
    from PyPDF2 import PdfReader
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        t = page.extract_text()
        if t and t.strip():
            texts.append(t.strip())
    return "\n".join(texts)


def parse_pptx(file_path):
    from pptx import Presentation
    prs = Presentation(file_path)
    texts = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
            if shape.has_notes_slide and shape.notes_slide.notes_text_frame:
                nt = shape.notes_slide.notes_text_frame.text.strip()
                if nt:
                    texts.append(f"[备注] {nt}")
    return "\n".join(texts)


def parse_xlsx(file_path):
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            rows.append(" | ".join(cells))
        lines.append(f"[工作表: {sheet_name}]\n" + "\n".join(rows))
    return "\n\n".join(lines)


def parse_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def parse(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".docx":
        return parse_docx(file_path)
    elif ext == ".pdf":
        return parse_pdf(file_path)
    elif ext == ".pptx":
        return parse_pptx(file_path)
    elif ext == ".xlsx":
        return parse_xlsx(file_path)
    elif ext == ".txt":
        return parse_txt(file_path)
    else:
        raise ValueError(f"不支持的文档格式: {ext}")
